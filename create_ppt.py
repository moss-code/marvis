#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
创建 2025年下半年营业厅经营分析报告.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 颜色方案 (Ocean Gradient风格)
COLOR_DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)
COLOR_TEAL = RGBColor(0x1C, 0x72, 0x93)
COLOR_MIDNIGHT = RGBColor(0x21, 0x29, 0x5C)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_OFF_WHITE = RGBColor(0xF2, 0xF5, 0xF7)
COLOR_LIGHT_BG = RGBColor(0xE8, 0xF0, 0xF5)
COLOR_ACCENT = RGBColor(0x02, 0x80, 0x90)
COLOR_GOLD = RGBColor(0xF9, 0xE7, 0x95)
COLOR_DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_GRAY_TEXT = RGBColor(0x5D, 0x6D, 0x7E)
COLOR_LIGHT_GRAY = RGBColor(0xBD, 0xC3, 0xCB)
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)
COLOR_GREEN = RGBColor(0x27, 0xAE, 0x60)
COLOR_TABLE_HEADER = RGBColor(0x06, 0x5A, 0x82)
COLOR_TABLE_ALT = RGBColor(0xE8, 0xF0, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='微软雅黑',
                anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines_data, font_name='微软雅黑'):
    """lines_data: list of (text, font_size, bold, color, alignment, space_after)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines_data):
        text, font_size, bold, color, alignment = item[:5]
        space_after = item[5] if len(item) > 5 else Pt(6)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = space_after
    return txBox


# ============================================================
# 第1页 - 封面页
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 背景 - 深蓝渐变
bg_shape = add_shape(slide1, 0, 0, W, H, fill_color=COLOR_DEEP_BLUE)

# 左侧装饰条
add_shape(slide1, 0, 0, Inches(0.4), H, fill_color=COLOR_GOLD)

# 主标题
add_textbox(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
            "2025年下半年营业厅经营分析报告",
            font_size=44, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.LEFT)

# 装饰线
add_shape(slide1, Inches(1.5), Inches(3.7), Inches(3), Inches(0.06), fill_color=COLOR_GOLD)

# 副标题
add_textbox(slide1, Inches(1.5), Inches(4.0), Inches(6), Inches(0.8),
            "2025年7月  |  数据驱动的业务决策",
            font_size=22, bold=False, color=RGBColor(0xBD, 0xE0, 0xEE), alignment=PP_ALIGN.LEFT)

# 底部信息
add_textbox(slide1, Inches(1.5), Inches(6.2), Inches(6), Inches(0.5),
            "中国电信 · 营业厅运营管理部",
            font_size=14, bold=False, color=RGBColor(0x8B, 0xBC, 0xCF), alignment=PP_ALIGN.LEFT)


# ============================================================
# 第2页 - 目录
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# 浅色背景
add_shape(slide2, 0, 0, W, H, fill_color=COLOR_OFF_WHITE)

# 顶部装饰条
add_shape(slide2, 0, 0, W, Inches(0.08), fill_color=COLOR_DEEP_BLUE)

# 标题
add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(5), Inches(0.8),
            "目  录", font_size=36, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)
add_textbox(slide2, Inches(0.8), Inches(1.15), Inches(5), Inches(0.4),
            "CONTENTS", font_size=12, bold=False, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.LEFT)

# 四个目录项
toc_items = [
    ("01", "业务办理概况", "各营业厅业务办理数据总览"),
    ("02", "用户增长分析", "净增用户数据与趋势"),
    ("03", "投诉分析", "客户投诉数据与问题定位"),
    ("04", "总结与建议", "核心结论与行动建议"),
]

colors_circle = [
    RGBColor(0x06, 0x5A, 0x82),
    RGBColor(0x1C, 0x72, 0x93),
    RGBColor(0x02, 0x80, 0x90),
    RGBColor(0x21, 0x29, 0x5C),
]

for i, (num, title, desc) in enumerate(toc_items):
    x = Inches(0.8) + (i % 4) * Inches(3.1)
    y = Inches(2.0)

    # 数字圆
    circle = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.8), y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors_circle[i]
    circle.line.fill.background()
    # 数字文字
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = '微软雅黑'
    p.alignment = PP_ALIGN.CENTER

    # 标题
    add_textbox(slide2, x, y + Inches(1.2), Inches(2.8), Inches(0.5),
                title, font_size=20, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.CENTER)

    # 描述
    add_textbox(slide2, x, y + Inches(1.7), Inches(2.8), Inches(0.5),
                desc, font_size=12, bold=False, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第3页 - 核心数据概览（表格）
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# 浅色背景
add_shape(slide3, 0, 0, W, H, fill_color=COLOR_OFF_WHITE)

# 顶部装饰条
add_shape(slide3, 0, 0, W, Inches(0.08), fill_color=COLOR_DEEP_BLUE)

add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "核心数据概览", font_size=32, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)

add_textbox(slide3, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "2025年下半年各营业厅业务办理数据（单位：户）", font_size=14, bold=False, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.LEFT)

# 表格数据
headers = ["营业厅", "宽带新装", "套餐办理", "5G升级", "净增用户"]
data_rows = [
    ["城东厅", "466", "798", "230", "269"],
    ["高新区厅", "396", "632", "176", "228"],
    ["城南厅", "338", "538", "163", "173"],
    ["城西厅", "304", "508", "136", "201"],
    ["城北厅", "274", "509", "132", "155"],
]

rows = len(data_rows) + 1
cols = len(headers)

left = Inches(0.8)
top = Inches(1.6)
width = Inches(11.7)
height = Inches(0.5) * rows

table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# 列宽
col_widths = [Inches(2.0), Inches(2.4), Inches(2.4), Inches(2.4), Inches(2.5)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# 表头
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 数据行
for r_idx, row_data in enumerate(data_rows):
    for c_idx, val in enumerate(row_data):
        cell = table.cell(r_idx + 1, c_idx)
        cell.text = val
        if r_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_TABLE_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.color.rgb = COLOR_DARK_TEXT
            p.font.name = '微软雅黑'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# 高亮列 - 净增用户列加粗
for r_idx in range(len(data_rows)):
    cell = table.cell(r_idx + 1, 4)
    for p in cell.text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

# 底部说明
add_textbox(slide3, Inches(0.8), Inches(5.8), Inches(10), Inches(0.4),
            "数据来源：2025年7月营业厅业务系统统计", font_size=10, bold=False,
            color=COLOR_LIGHT_GRAY, alignment=PP_ALIGN.LEFT)


# ============================================================
# 第4页 - 业务排名
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide4, 0, 0, W, H, fill_color=COLOR_OFF_WHITE)
add_shape(slide4, 0, 0, W, Inches(0.08), fill_color=COLOR_DEEP_BLUE)

add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "业务办理排名", font_size=32, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)

# ---- 宽带新装排名 ----
add_textbox(slide4, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
            "宽带新装排名", font_size=20, bold=True, color=COLOR_TEAL, alignment=PP_ALIGN.LEFT)

# 条形图（模拟）
biz_data_1 = [("城东厅", 466), ("高新区厅", 396), ("城南厅", 338), ("城西厅", 304), ("城北厅", 274)]
max_val = max(v for _, v in biz_data_1)

bar_colors = [
    RGBColor(0xF9, 0xE7, 0x95),  # 金色
    RGBColor(0xBD, 0xC3, 0xCB),  # 灰色
    RGBColor(0xCD, 0x7F, 0x32),  # 铜色
    COLOR_TEAL,
    COLOR_LIGHT_GRAY,
]

for i, (name, val) in enumerate(biz_data_1):
    y = Inches(1.9) + i * Inches(0.7)
    bar_w = int(Inches(6.0) * val / max_val) if max_val > 0 else Inches(0.1)
    # 名称
    add_textbox(slide4, Inches(0.8), y, Inches(1.4), Inches(0.5),
                name, font_size=13, bold=True, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    # 条
    bar_shape = add_shape(slide4, Inches(2.3), y + Inches(0.05), bar_w, Inches(0.4),
                          fill_color=bar_colors[i] if i < 3 else COLOR_TEAL)
    # 数值
    add_textbox(slide4, Inches(2.3) + bar_w + Inches(0.1), y + Inches(0.05),
                Inches(1.5), Inches(0.4),
                str(val), font_size=13, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)

# ---- 套餐办理排名 ----
add_textbox(slide4, Inches(6.8), Inches(1.3), Inches(5), Inches(0.5),
            "套餐办理排名", font_size=20, bold=True, color=COLOR_TEAL, alignment=PP_ALIGN.LEFT)

biz_data_2 = [("城东厅", 798), ("高新区厅", 632), ("城南厅", 538), ("城北厅", 509), ("城西厅", 508)]
max_val2 = max(v for _, v in biz_data_2)

for i, (name, val) in enumerate(biz_data_2):
    y = Inches(1.9) + i * Inches(0.7)
    bar_w = int(Inches(6.0) * val / max_val2) if max_val2 > 0 else Inches(0.1)
    add_textbox(slide4, Inches(6.8), y, Inches(1.4), Inches(0.5),
                name, font_size=13, bold=True, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    bar_shape = add_shape(slide4, Inches(8.3), y + Inches(0.05), bar_w, Inches(0.4),
                          fill_color=bar_colors[i] if i < 3 else COLOR_TEAL)
    add_textbox(slide4, Inches(8.3) + bar_w + Inches(0.1), y + Inches(0.05),
                Inches(1.5), Inches(0.4),
                str(val), font_size=13, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)

# 图例
add_textbox(slide4, Inches(0.8), Inches(6.5), Inches(12), Inches(0.3),
            "🥇 第一名  🥈 第二名  🥉 第三名", font_size=11, bold=False,
            color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.LEFT)


# ============================================================
# 第5页 - 用户增长分析
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide5, 0, 0, W, H, fill_color=COLOR_OFF_WHITE)
add_shape(slide5, 0, 0, W, Inches(0.08), fill_color=COLOR_DEEP_BLUE)

add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "用户增长分析", font_size=32, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)
add_textbox(slide5, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "各营业厅净增用户数据", font_size=14, bold=False, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.LEFT)

# 净增用户排名 - 大数字展示
growth_data = [("城东厅", 269, True), ("高新区厅", 228, True), ("城西厅", 201, True),
               ("城南厅", 173, True), ("城北厅", 155, True)]

for i, (name, val, positive) in enumerate(growth_data):
    x = Inches(0.8) + (i % 5) * Inches(2.45)
    y = Inches(1.6)

    # 卡片背景
    card = add_shape(slide5, x, y, Inches(2.2), Inches(2.0),
                     fill_color=COLOR_WHITE)
    card.shadow.inherit = False

    # 大数字
    color_num = COLOR_GREEN if positive else COLOR_RED
    prefix = "+" if positive else ""
    add_textbox(slide5, x + Inches(0.1), y + Inches(0.3), Inches(2.0), Inches(1.0),
                f"{prefix}{val}", font_size=42, bold=True, color=color_num, alignment=PP_ALIGN.CENTER)

    # 营业厅名
    add_textbox(slide5, x + Inches(0.1), y + Inches(1.3), Inches(2.0), Inches(0.5),
                name, font_size=16, bold=True, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.CENTER)

# 底部说明框
note_shape = add_shape(slide5, Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.5),
                       fill_color=RGBColor(0xD5, 0xEB, 0xEE))
add_textbox(slide5, Inches(1.2), Inches(4.4), Inches(11), Inches(1.2),
            "📈  关键发现\n\n"
            "• 所有营业厅均为正增长，用户留存状况良好\n"
            "• 城东厅净增 269 户领跑，高新区厅 228 户紧随其后\n"
            "• 城北厅净增 155 户排名靠后，建议加大营销力度",
            font_size=14, bold=False, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)


# ============================================================
# 第6页 - 投诉分析
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide6, 0, 0, W, H, fill_color=COLOR_OFF_WHITE)
add_shape(slide6, 0, 0, W, Inches(0.08), fill_color=COLOR_DEEP_BLUE)

add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
            "投诉分析", font_size=32, bold=True, color=COLOR_DEEP_BLUE, alignment=PP_ALIGN.LEFT)
add_textbox(slide6, Inches(0.8), Inches(1.0), Inches(8), Inches(0.4),
            "各营业厅客户投诉数据", font_size=14, bold=False, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.LEFT)

# 投诉排名 - 水平条（反向，投诉多的在前面）
complaint_data = [("城东厅", 46), ("城南厅", 37), ("高新区厅", 30), ("城西厅", 17), ("城北厅", 15)]
max_c = max(v for _, v in complaint_data)

add_textbox(slide6, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
            "投诉总量排名（需重点关注）", font_size=18, bold=True, color=COLOR_RED, alignment=PP_ALIGN.LEFT)

for i, (name, val) in enumerate(complaint_data):
    y = Inches(2.2) + i * Inches(0.7)
    bar_w = int(Inches(5.0) * val / max_c) if max_c > 0 else Inches(0.1)
    add_textbox(slide6, Inches(0.8), y, Inches(1.4), Inches(0.5),
                name, font_size=13, bold=True, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.RIGHT)

    # 颜色随投诉量变化
    if val >= 40:
        bar_color = COLOR_RED
    elif val >= 25:
        bar_color = RGBColor(0xE6, 0x7E, 0x22)
    else:
        bar_color = COLOR_GREEN

    bar_shape = add_shape(slide6, Inches(2.3), y + Inches(0.05), bar_w, Inches(0.4),
                          fill_color=bar_color)
    add_textbox(slide6, Inches(2.3) + bar_w + Inches(0.1), y + Inches(0.05),
                Inches(1.5), Inches(0.4),
                str(val) + " 件", font_size=13, bold=True, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)

# 右侧重点提示
right_box = add_shape(slide6, Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.5),
                      fill_color=COLOR_WHITE)
add_textbox(slide6, Inches(7.8), Inches(1.7), Inches(4.4), Inches(0.5),
            "⚠️ 重点关注", font_size=18, bold=True, color=COLOR_RED, alignment=PP_ALIGN.LEFT)

add_textbox(slide6, Inches(7.8), Inches(2.3), Inches(4.4), Inches(3.0),
            "1. 城东厅投诉量 46 件，为全厅最高\n"
            "   建议：加强服务品质管控\n\n"
            "2. 网络质量类投诉需重点关注\n"
            "   建议：排查网络覆盖问题\n\n"
            "3. 城北厅投诉量最低 15 件\n"
            "   值得其他营业厅借鉴经验",
            font_size=14, bold=False, color=COLOR_DARK_TEXT, alignment=PP_ALIGN.LEFT)


# ============================================================
# 第7页 - 总结与建议
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])

# 深色背景
add_shape(slide7, 0, 0, W, H, fill_color=COLOR_DEEP_BLUE)

# 标题
add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "总结与建议", font_size=36, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.LEFT)

add_shape(slide7, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.05), fill_color=COLOR_GOLD)

# 四个建议卡片
suggestions = [
    ("🏆", "城东厅", "业务量最大，但投诉也最多\n需加强服务质量管控"),
    ("⭐", "高新区厅", "业务和用户增长均表现优秀\n可作为标杆示范厅"),
    ("📢", "城北厅", "业务量偏低\n可考虑加大营销推广力度"),
    ("🔧", "网络质量", "各厅网络质量类投诉\n需重点关注和排查"),
]

card_colors = [
    RGBColor(0x08, 0x7A, 0x9F),
    RGBColor(0x0A, 0x8C, 0xB0),
    RGBColor(0x0C, 0x9E, 0xC2),
    RGBColor(0x0E, 0xB0, 0xD4),
]

for i, (icon, title, desc) in enumerate(suggestions):
    x = Inches(0.8) + (i % 4) * Inches(3.1)
    y = Inches(2.0)

    card = add_shape(slide7, x, y, Inches(2.8), Inches(4.2), fill_color=card_colors[i])

    # 图标
    add_textbox(slide7, x, y + Inches(0.3), Inches(2.8), Inches(0.8),
                icon, font_size=36, bold=False, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    # 标题
    add_textbox(slide7, x, y + Inches(1.2), Inches(2.8), Inches(0.5),
                title, font_size=20, bold=True, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

    # 装饰线
    add_shape(slide7, x + Inches(0.6), y + Inches(1.8), Inches(1.6), Inches(0.03),
              fill_color=COLOR_GOLD)

    # 描述
    add_textbox(slide7, x + Inches(0.2), y + Inches(2.0), Inches(2.4), Inches(2.0),
                desc, font_size=13, bold=False, color=RGBColor(0xD5, 0xEB, 0xEE),
                alignment=PP_ALIGN.CENTER)


# ============================================================
# 保存文件
# ============================================================
output_path = r"D:\study\电信\实操\培训\marvis\2025年下半年营业厅经营分析报告.pptx"
prs.save(output_path)
file_size = os.path.getsize(output_path)
print(f"文件已生成: {output_path}")
print(f"文件大小: {file_size / 1024:.1f} KB")
